#!/usr/bin/env python3
"""Generate DeOpenRouter course presentation (PPTX) with speaker notes."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


SLIDES = [
    {
        "title": "DeOpenRouter: A Low-Trust AI API Marketplace Prototype",
        "bullets": [
            "Solidity marketplace + Next.js + Hono relay + optional FastAPI audit",
            "Hybrid: inference off-chain; coordination, pricing, payment, audit anchors on-chain",
        ],
        "notes": """Good morning. I'm presenting DeOpenRouter, a course project that prototypes an AI API marketplace with a smaller trust footprint than a typical centralized gateway. It is built with a Solidity marketplace contract, a Next.js frontend, a Hono relay API, and an optional FastAPI audit service. The headline idea is hybrid: keep model inference off-chain for practicality, while moving coordination, pricing, payments, and audit anchoring to an EVM chain.""",
    },
    {
        "title": "The Problem: Why Centralized AI Gateways Are “Trust-Heavy”",
        "bullets": [
            "One operator often controls pricing, routing, billing, and disputes",
            "Hard to verify quiet rule changes, fair charges, or rewrite-resistant records",
            "Question: what can we make publicly observable without full on-chain inference?",
        ],
        "notes": """Today, most AI API routers ask users to trust one operator for pricing, routing, billing, and disputes. That is convenient, but it is hard for users to independently verify whether rules changed quietly, whether a charge matched the advertised policy, or whether records could be rewritten after the fact. Our story starts from that gap: what parts of the workflow can we make publicly observable and harder to tamper with, without pretending we can run full inference on-chain?""",
    },
    {
        "title": "Project Goal: Move Key Trust Boundaries On-Chain",
        "bullets": [
            "Public provider registration + commitments; delayed, observable price changes",
            "Auditable record per paid call; stake + slashing proposal with challenge window",
            "External audits anchorable as hashes; inference stays off-chain for realism",
        ],
        "notes": """DeOpenRouter tries to migrate a few critical trust boundaries to Ethereum-style execution: provider commitments and metadata are registered publicly, price changes are not instantly sneaky but have observable delay, each paid call leaves an auditable on-chain record, stake can be reduced through an on-chain slashing proposal with a challenge window, and external audit results can be anchored as hashes. Real inference still happens off-chain, which keeps the system realistic while improving transparency for settlement and auditing.""",
    },
    {
        "title": "Design Choice: Hybrid Architecture (On-Chain vs Off-Chain)",
        "bullets": [
            "On-chain: register, stake, delayed pricing, per-call settlement, slashing, audit anchors",
            "Off-chain: relay inference (local echo dev mode or OpenRouter proxy)",
            "Explicit trade-off: clearer settlement/records, not a full decentralized inference network",
        ],
        "notes": """We deliberately use a hybrid architecture. On-chain, the marketplace handles registration, staking, delayed pricing, per-call settlement, slashing governance, and audit anchoring. Off-chain, the relay performs inference—either a local echo mode for development or proxying to OpenRouter when configured. This split is the core design trade-off: we do not claim a fully decentralized inference network; we claim a clearer, more checkable settlement and recordkeeping layer implemented in Solidity.""",
    },
    {
        "title": "What’s in the Repo (MVP Scope)",
        "bullets": [
            "MVP focused on trust + settlement model (README), not production-scale decentralization",
            "Foundry contracts · Next.js (user + provider) · Hono relay · FastAPI audit server",
            "Docs: AUDIT_GOVERNANCE.md, TRUST_LAYERS.md for assumptions and future directions",
        ],
        "notes": """The repository is an MVP prototype focused on trust and settlement mechanics, not production-scale decentralization. It includes the Foundry smart-contract marketplace, a Next.js web app with separate user and provider flows, a Hono relay API for chat completions, and an audit server that can structurally probe the relay and summarize results. Documentation also points to governance and trust-layer notes for anyone evaluating security assumptions.""",
    },
    {
        "title": "End-to-End User Journey (Five Steps)",
        "bullets": [
            "1 Provider registers: stake, metadata URI, identity/capability hashes, endpoint commitment",
            "2 Web lists providers; user simulates calls via relay",
            "3 Client hashes request + response",
            "4 User pays → invoke → fee transfer + CallRecorded event",
            "5 Optional audit: probe relay → canonical JSON → anchor keccak256 on-chain",
        ],
        "notes": """Here is the story as a workflow. First, a provider registers on-chain with stake, metadata URI, identity and capability hashes, and an endpoint commitment. Second, the web app lists providers and users can simulate calls through the relay. Third, the client hashes the request and response. Fourth, the user pays and calls invoke, which transfers fees and emits a CallRecorded event. Fifth, if auditing is enabled, the system probes the relay, canonicalizes JSON, optionally publishes off-chain, then anchors keccak256 of canonical JSON on-chain.""",
    },
    {
        "title": "Solidity Core: DeOpenRouterMarketplace.sol",
        "bullets": [
            "contracts/src/DeOpenRouterMarketplace.sol — on-chain source of truth",
            "Registration, effective pricing, per-call payment, event-based traceability",
            "Course anchor: concrete functions + state transitions, not “crypto-themed UI only”",
        ],
        "notes": """The heart of the implementation is DeOpenRouterMarketplace.sol. This contract is the on-chain source of truth for provider registration, effective pricing logic, per-call payments, and event-based traceability. For grading purposes, this slide is the anchor: the project is not a website with crypto vibes; it is a concrete Solidity marketplace with explicit functions and state transitions that implement the economic and recordkeeping rules we just described.""",
    },
    {
        "title": "Key Functions: Registration, Metadata, Delayed Pricing",
        "bullets": [
            "register — stake + commitments (not raw endpoint URL on-chain)",
            "updateProviderMetadata — metadata URI / hashes",
            "announcePriceChange + getEffectivePrice — delayed activation, read-only price read",
        ],
        "notes": """Providers enter through register, which includes staking and commitments rather than publishing raw endpoint URLs in plaintext on-chain. They can update certain metadata fields via updateProviderMetadata. Pricing credibility comes from announcePriceChange combined with getEffectivePrice, which reads the currently effective price without mutating state. The README emphasizes delayed activation: operators cannot silently flip prices in a way that users cannot anticipate if they watch the chain.""",
    },
    {
        "title": "Key Function: invoke (Pay, Hash, Record, Emit)",
        "bullets": [
            "Receives payment; stores request/response hashes (not full plaintext)",
            "Emits CallRecorded — verifiable on-chain footprint tied to provider + fee",
        ],
        "notes": """The user-facing settlement path is invoke. The contract receives payment, records cryptographic commitments to the request and response, and emits CallRecorded. This is intentionally not store the full prompt on-chain; it stores hashes for auditability and reconciliation while limiting on-chain data exposure. In presentation terms, this is the moment where AI usage becomes a verifiable on-chain event tied to a provider and a fee transfer.""",
    },
    {
        "title": "Governance, Slashing, and Audit Anchoring (MVP)",
        "bullets": [
            "Roles: propose/accept slash operator and audit recorder; slashTreasury + challenge period",
            "Slashing: proposeSlash → challenge window → finalizeSlashProposal",
            "Audit: recordAudit / recordAuditWithUri; README warns multisig/treasury assumptions",
        ],
        "notes": """Beyond the happy path, the contract includes governance hooks like proposing and accepting slash and audit recorder roles, setting a slash treasury and challenge period, and audit flows such as recordAudit, recordAuditWithUri, and multi-signer attestation patterns described in the docs. Slashing follows a proposal, optional challenge within a window, then finalization sending ETH to slashTreasury. The README warns economic safety still depends on multisig and treasury configuration—this is not a complete arbitration market.""",
    },
    {
        "title": "Honest MVP Boundaries (What We Are NOT Claiming)",
        "bullets": [
            "Hashes on-chain, not full prompts/responses; instant per-call settlement, not escrow",
            "No decentralized discovery/reputation; auditor quorum aggregation largely off-chain",
            "No cryptographic proof of correct inference — see TRUST_LAYERS.md for directions",
        ],
        "notes": """To score well on story, we should be explicit about limits. The chain stores commitments and hashes, not full plaintext prompts and responses. Settlement is instant per call, not escrowed streaming settlement. Endpoint discovery is not a decentralized reputation network. Multi-auditor quorum aggregation happens off-chain with on-chain events and anchors. There is no cryptographic proof of correct inference execution—those are future trust layers, referenced in the documentation rather than implemented here.""",
    },
    {
        "title": "Demo Path, Tests, and Takeaway",
        "bullets": [
            "Anvil → forge deploy → relay + web; wallet on chain id 31337; Provider then User flow",
            "Tests: forge test, web/api/audit-server test commands in README",
            "Takeaway: Solidity can enforce pricing discipline, payment, and tamper-resistant records around API usage while inference stays off-chain",
        ],
        "notes": """Locally, you run Anvil, deploy with Foundry, start the relay and web app, connect a wallet on chain id 31337, register as a provider, then invoke as a user and inspect logs or events. Tests exist for contracts, web, relay, and audit server as described in the README. My takeaway is that Solidity can credibly enforce pricing discipline, payment, and tamper-resistant records around AI API usage, while inference remains off-chain—useful as a classroom prototype and a foundation for stricter trust models later. Thank you—questions?""",
    },
]


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1].has_text_frame:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes


def main() -> None:
    out_dir = Path(__file__).resolve().parent
    out_path = out_dir / "DeOpenRouter-Course-Presentation.pptx"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "DeOpenRouter",
        "Decentralization & Cryptocurrency — Course Project\n"
        "Story · Solidity Implementation · Presentation",
    )
    # Speaker notes on title slide
    title_slide = prs.slides[0]
    title_slide.notes_slide.notes_text_frame.text = (
        "Introduce yourself and the course context. Mention grading dimensions briefly: "
        "this deck emphasizes story, then Solidity implementation, then a clear demo path."
    )

    for item in SLIDES:
        add_bullet_slide(prs, item["title"], item["bullets"], item["notes"])

    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
